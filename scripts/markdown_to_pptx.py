#!/usr/bin/env python3
import argparse
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


IMAGE_PATTERN = re.compile(r"!\[[^\]]*\]\(([^\)]+)\)")


@dataclass
class Subcategory:
    name: str
    bullets: List[str] = field(default_factory=list)


@dataclass
class MainSection:
    name: str
    subcategories_chart: Optional[Path] = None
    subcategories: List[Subcategory] = field(default_factory=list)


@dataclass
class MarkdownReport:
    title: str
    main_chart: Optional[Path]
    sections: List[MainSection]


def parse_markdown(md_path: Path) -> MarkdownReport:
    text = md_path.read_text(encoding="utf-8")
    lines = [line.rstrip("\n") for line in text.splitlines()]

    title = ""
    main_chart: Optional[Path] = None
    sections: List[MainSection] = []

    i = 0
    base_dir = md_path.parent

    # Tytuł: pierwsza linia zaczynająca się od '# '
    while i < len(lines):
        line = lines[i].strip()
        if line.startswith("# "):
            title = line[2:].strip()
            i += 1
            break
        i += 1

    # Po tytule spodziewamy się obrazu kategorii głównych
    while i < len(lines):
        line = lines[i].strip()
        m = IMAGE_PATTERN.search(line)
        if m:
            img_rel = m.group(1).strip()
            main_chart = (base_dir / img_rel).resolve()
            i += 1
            break
        if line.startswith("## "):
            # brak obrazu main chart
            break
        i += 1

    current_section: Optional[MainSection] = None
    current_sub: Optional[Subcategory] = None

    while i < len(lines):
        raw = lines[i]
        line = raw.strip()
        if line.startswith("## "):
            # Zapisz poprzedni podrozdział, jeśli był
            if current_sub and current_section:
                current_section.subcategories.append(current_sub)
            # Zapisz poprzednią sekcję, jeśli była
            if current_section:
                sections.append(current_section)
            current_section = MainSection(name=line[3:].strip())
            current_sub = None
            # Sprawdź, czy kolejna linia to obraz podkategorii
            j = i + 1
            while j < len(lines):
                next_line = lines[j].strip()
                if next_line == "":
                    j += 1
                    continue
                m = IMAGE_PATTERN.search(next_line)
                if m:
                    img_rel = m.group(1).strip()
                    current_section.subcategories_chart = (base_dir / img_rel).resolve()
                    break
                if next_line.startswith("### ") or next_line.startswith("## "):
                    break
                j += 1
            i += 1
            continue
        elif line.startswith("### "):
            # Zapisz poprzedni sub, jeśli był
            if current_sub and current_section:
                current_section.subcategories.append(current_sub)
            current_sub = Subcategory(name=line[4:].strip())
            i += 1
            continue
        elif line.startswith("- ") and current_sub is not None:
            bullet = line[2:].strip()
            if bullet:
                current_sub.bullets.append(bullet)
            i += 1
            continue
        else:
            i += 1

    # Zapisz końcowe elementy
    if current_sub and current_section:
        current_section.subcategories.append(current_sub)
    if current_section:
        sections.append(current_section)

    return MarkdownReport(title=title or md_path.stem, main_chart=main_chart, sections=sections)


def set_presentation_size_a4_landscape(prs: Presentation) -> None:
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)


def style_title(shape, font_size: int = 40) -> None:
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    p.alignment = PP_ALIGN.LEFT


def style_body_paragraph(p, size_pt: int = 16, level: int = 0) -> None:
    p.level = level
    p.font.name = "Calibri"
    p.font.size = Pt(size_pt)
    p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    # ciaśniejsze interlinie dla czytelności i oszczędności miejsca
    try:
        p.line_spacing = 1.0
        p.space_after = Pt(2)
        p.space_before = Pt(0)
    except Exception:
        pass


def add_title_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    # Dodaj pole tytułu manualnie, aby kontrolować pozycję
    left = Inches(0.8)
    top = Inches(0.8)
    width = Inches(10.0)
    height = Inches(1.5)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    title_shape.text = title
    style_title(title_shape, font_size=40)


def add_image_slide(prs: Presentation, title: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Tytuł
    title_shape = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10.0), Inches(1.0))
    title_shape.text = title
    style_title(title_shape, font_size=26)

    # Obraz – dopasowanie do obszaru z marginesami
    if not image_path.exists():
        warn = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(9.0), Inches(1.0))
        warn.text = f"Brak obrazu: {image_path.name}"
        style_body_paragraph(warn.text_frame.paragraphs[0], size_pt=16)
        return

    max_left = Inches(0.8)
    max_top = Inches(1.6)
    max_width = Inches(10.0)
    max_height = Inches(6.0)

    pic = slide.shapes.add_picture(str(image_path), max_left, max_top)

    # Skalowanie aby zmieścić całość w ramce
    scale_w = max_width / pic.width
    scale_h = max_height / pic.height
    scale = min(scale_w, scale_h)

    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    # Wycentruj w ramce poziomo
    pic.left = int(max_left + (max_width - pic.width) / 2)
    pic.top = int(max_top + (max_height - pic.height) / 2)


def add_bullets_slide(prs: Presentation, title: str, bullets: List[str]) -> None:
    # Blank layout + własny szeroki textbox dla większej czytelności
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Tytuł
    title_shape = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10.0), Inches(0.9))
    title_shape.text = title
    style_title(title_shape, font_size=24)

    # Główne pole tekstowe na treść – szerokie, z małymi marginesami i zawijaniem
    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(10.0)
    height = Inches(5.9)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    try:
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.word_wrap = True
    except Exception:
        pass

    max_items = 12  # limit elementów na slajd dla przejrzystości
    chunk = bullets[:max_items]

    for idx, item in enumerate(chunk):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # manualny punktor, aby nie tracić kontroli nad formatowaniem
        p.text = f"• {item}"
        style_body_paragraph(p, size_pt=16, level=0)


def build_pptx(md_path: str, output_path: str) -> str:
    md = parse_markdown(Path(md_path))

    prs = Presentation()
    set_presentation_size_a4_landscape(prs)

    # Slajd tytułowy
    add_title_slide(prs, md.title)

    # Slajd z głównym wykresem
    if md.main_chart is not None:
        add_image_slide(prs, "Kategorie główne (malejąco)", md.main_chart)

    # Sekcje
    for section in md.sections:
        # Slajd z wykresem podkategorii
        if section.subcategories_chart is not None:
            add_image_slide(prs, f"{section.name} – podkategorie (malejąco)", section.subcategories_chart)
        else:
            add_image_slide(prs, f"{section.name} – podkategorie (malejąco)", Path("(brak obrazu)"))
        # Slajdy z przykładami
        for sub in section.subcategories:
            title = f"{section.name} / {sub.name}"
            bullets = sub.bullets[:10]
            if not bullets:
                bullets = ["(Brak streszczeń)"]
            add_bullets_slide(prs, title, bullets)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Konwertuje raport Markdown (z obrazami PNG) do prezentacji PPTX")
    parser.add_argument("--markdown", required=True, help="Ścieżka do pliku .md")
    parser.add_argument("--output", required=False, help="Ścieżka wyjściowa .pptx")
    return parser.parse_args()


essential_packages = ["python-pptx", "Pillow"]

def ensure_packages_installed() -> None:
    # Ślepe importy sprawdzą dostępność; jeśli brak, zasygnalizujemy w CLI (nie instalujemy automatycznie)
    try:
        import pptx  # noqa: F401
    except Exception as exc:
        raise SystemExit("Brak pakietu 'python-pptx'. Zainstaluj: pip install python-pptx Pillow") from exc


def main() -> int:
    ensure_packages_installed()
    args = parse_args()
    md_path = os.path.abspath(args.markdown)
    if not os.path.exists(md_path):
        raise FileNotFoundError(f"Nie znaleziono pliku Markdown: {md_path}")

    default_out = os.path.splitext(md_path)[0] + ".pptx"
    output_path = os.path.abspath(args.output) if args.output else default_out

    result = build_pptx(md_path=md_path, output_path=output_path)
    print(result)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
